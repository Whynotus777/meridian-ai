"""IC Deck (PowerPoint) export for CIM analysis results.

Generates an 11-slide investment committee presentation from an AnalysisResult.

Slide structure:
  1  Title slide
  2  Deal Overview         (hero metrics + grade)
  3  Executive Summary
  4  Company Overview
  5  Financial Summary     (two-column income/balance)
  6  Revenue Analysis      (bar chart + segments)
  7  Growth Thesis
  8  Key Risks & Mitigants (color-coded table)
  9  Comparable Companies
  10 Deal Scoring          (dimension bars)
  11 Key Diligence Questions + Recommendation

Requires: python-pptx >= 1.0
"""

from __future__ import annotations

import json
import os
import re
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple, TYPE_CHECKING

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from core.pipeline import AnalysisResult


# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
_NAVY       = RGBColor(0x1E, 0x27, 0x61)   # primary dark bg
_NAVY_MID   = RGBColor(0x2B, 0x3A, 0x8A)   # header bar
_ICE_BLUE   = RGBColor(0xCA, 0xDC, 0xFC)   # text on dark slides
_ORANGE     = RGBColor(0xE8, 0x73, 0x4A)   # Meridian accent
_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
_OFF_WHITE  = RGBColor(0xF5, 0xF5, 0xF5)   # content slide bg
_LIGHT_GREY = RGBColor(0xE8, 0xE8, 0xE8)   # alt table rows
_MED_GREY   = RGBColor(0xBB, 0xBB, 0xBB)   # borders / captions
_DARK_GREY  = RGBColor(0x44, 0x44, 0x44)   # body text
_BLACK      = RGBColor(0x1A, 0x1A, 0x1A)   # headings
_GREEN      = RGBColor(0x4C, 0xAF, 0x50)
_AMBER      = RGBColor(0xFF, 0xC1, 0x07)
_RED        = RGBColor(0xEF, 0x53, 0x50)
_DARK_RED   = RGBColor(0xC0, 0x00, 0x00)

_GRADE_COLORS: Dict[str, RGBColor] = {
    "A": RGBColor(0x2E, 0x7D, 0x32),
    "B": RGBColor(0x43, 0xA0, 0x47),
    "C": _AMBER,
    "D": _RED,
    "F": _DARK_RED,
}
_SEV_COLORS: Dict[str, RGBColor] = {
    "Critical": _DARK_RED,
    "High":     _RED,
    "Medium":   _AMBER,
    "Low":      _GREEN,
}

# Slide canvas — widescreen 16:9
_W      = Inches(13.33)
_H      = Inches(7.5)
_MARGIN = Inches(0.5)
_FOOTER_H = Inches(0.30)


# ---------------------------------------------------------------------------
# Presentation + slide scaffolding
# ---------------------------------------------------------------------------

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------

def _rect(slide, l, t, w, h, fill: RGBColor, border: Optional[RGBColor] = None):
    """Add a solid-filled rectangle. No border by default."""
    shp = slide.shapes.add_shape(1, l, t, w, h)   # 1 = RECTANGLE
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()          # no border
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.5)
    return shp


def _label(slide, text: str, l, t, w, h, *,
           bold=False, size=12, color=_DARK_GREY,
           align=PP_ALIGN.LEFT, italic=False, wrap=True) -> None:
    """Add a plain text box with a single paragraph."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text      = str(text)
    p.alignment = align
    for run in p.runs:
        run.font.name   = "Calibri"
        run.font.bold   = bold
        run.font.italic = italic
        run.font.size   = Pt(size)
        run.font.color.rgb = color


def _multiline(slide, lines: List[Tuple[str, dict]], l, t, w, h) -> None:
    """Add a text box with multiple paragraphs (list of (text, kwargs))."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (text, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text      = str(text)
        p.alignment = kw.get("align", PP_ALIGN.LEFT)
        sp = kw.get("space_before", 0)
        if sp:
            p.space_before = Pt(sp)
        for run in p.runs:
            run.font.name    = "Calibri"
            run.font.bold    = kw.get("bold", False)
            run.font.italic  = kw.get("italic", False)
            run.font.size    = Pt(kw.get("size", 13))
            run.font.color.rgb = kw.get("color", _DARK_GREY)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Slide chrome — header bar + footer
# ---------------------------------------------------------------------------

def _header(slide, title: str) -> None:
    """Navy title bar across the top of a content slide."""
    _rect(slide, 0, 0, _W, Inches(0.68), _NAVY)
    _label(slide, title,
           _MARGIN, Inches(0.10), _W - _MARGIN * 2, Inches(0.50),
           bold=True, size=22, color=_WHITE)


def _footer(slide, slide_num: int) -> None:
    """'Meridian AI · PE Deal Intelligence' footer + slide number."""
    fy = _H - _FOOTER_H - Inches(0.04)
    _rect(slide, 0, fy, _W, _FOOTER_H, _NAVY)
    _label(slide, "Meridian AI  ·  PE Deal Intelligence",
           _MARGIN, fy + Inches(0.04), Inches(9), _FOOTER_H - Inches(0.04),
           size=8.5, color=_MED_GREY)
    _label(slide, str(slide_num),
           _W - Inches(0.65), fy + Inches(0.04), Inches(0.55), _FOOTER_H - Inches(0.04),
           size=8.5, color=_MED_GREY, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Data formatting helpers
# ---------------------------------------------------------------------------

def _fmt_num(val) -> str:
    if val is None or str(val).lower() in ("not_provided", "none", ""):
        return "N/A"
    try:
        v = float(val)
    except (ValueError, TypeError):
        return str(val)
    sign = "-" if v < 0 else ""
    av = abs(v)
    if av >= 1e9:  return f"{sign}${av/1e9:.1f}B"
    if av >= 1e6:  return f"{sign}${av/1e6:.1f}M"
    if av >= 1e3:  return f"{sign}${av/1e3:.0f}K"
    return f"{sign}${av:,.0f}"


def _fmt_pct(val) -> str:
    if val is None or str(val).lower() in ("not_provided", "none", ""):
        return "N/A"
    try:
        return f"{float(val):.1%}"
    except (ValueError, TypeError):
        return str(val)


def _fmt_str(val) -> str:
    if val is None or str(val).lower() in ("not_provided", "none", ""):
        return "N/A"
    return str(val)


def _with_currency(num_str: str, currency: str) -> str:
    """Replace '$' prefix with currency code for non-USD docs."""
    if currency in ("USD", "N/A", ""):
        return num_str
    return num_str.replace("$", f"{currency} ")


def _clean_markdown(text: str) -> str:
    """Strip markdown formatting for clean PPTX rendering."""
    if not text:
        return text
    # Bold / italic markers
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text, flags=re.DOTALL)
    text = re.sub(r'\*(.+?)\*',     r'\1', text, flags=re.DOTALL)
    # Heading markers
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    # Horizontal rules
    text = re.sub(r'^[-_]{3,}\s*$', '', text, flags=re.MULTILINE)
    # Leading bullet markers (* - • ·) at start of lines
    text = re.sub(r'^\s*[\*\-•·]\s+', '', text, flags=re.MULTILINE)
    # Inline code backticks
    text = re.sub(r'`+', '', text)
    # Collapse excessive blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _truncate_at_word(text: str, max_chars: int = 150) -> str:
    """Truncate text at a word boundary; append '…' if truncated."""
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars]
    last_space = cut.rfind(' ')
    if last_space > max_chars * 0.55:
        cut = cut[:last_space]
    return cut.rstrip(' .,;') + "…"


def _parse_numbered_items(content: str, max_items: int = 8) -> List[str]:
    """Extract only explicitly numbered or bulleted items; skip preamble prose."""
    if not content:
        return []
    content = _clean_markdown(content)
    items: List[str] = []
    for line in content.split("\n"):
        line = line.strip()
        m = re.match(r'^(?:\d+[\.\)]\s+|\s*[\*\-•·]\s+)(.+)', line)
        if m:
            item = m.group(1).strip()
            if len(item) > 15:
                items.append(item)
                if len(items) >= max_items:
                    break
    return items


def _cite_page(obj: Any) -> Optional[int]:
    if not isinstance(obj, dict):
        return None
    pg = obj.get("page")
    if isinstance(pg, int) and pg > 0:
        return pg
    return None


# ---------------------------------------------------------------------------
# Memo parsing helpers
# ---------------------------------------------------------------------------

def _parse_memo(memo_str: str) -> Dict[str, Any]:
    if not memo_str:
        return {}
    try:
        return json.loads(memo_str)
    except (json.JSONDecodeError, ValueError):
        m = re.search(r'\{.*\}', memo_str, re.DOTALL)
        if m:
            try:
                return json.loads(m.group(0))
            except Exception:
                pass
    return {}


def _find_section(sections: List[Dict], *keywords) -> str:
    """Return content from first section whose heading matches any keyword."""
    kws = [k.lower() for k in keywords]
    for s in sections:
        heading = s.get("heading", s.get("title", "")).lower()
        if any(kw in heading for kw in kws):
            return s.get("content", "")
    return ""


def _to_bullets(content: str, max_bullets: int = 6) -> List[str]:
    """Extract short bullet-worthy lines from prose content."""
    if not content:
        return []
    content = _clean_markdown(content)
    bullets: List[str] = []
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        line = re.sub(r'^[\*\-•·\d]+[\.\)]\s*', '', line).strip()
        # Skip very short, very long, or intro-style lines
        if len(line) < 18 or len(line) > 260:
            continue
        if re.match(
            r'.*(following|include:|as follows|below|several risks|key risks|'
            r'inherent|outlined below|are highlighted|questions include)',
            line, re.I,
        ) and len(line) < 120:
            continue
        bullets.append(line)
        if len(bullets) >= max_bullets:
            break
    # Fall back to sentence splitting if no usable lines
    if not bullets:
        for sent in re.split(r'(?<=[.!?])\s+', content):
            sent = sent.strip()
            if 20 < len(sent) < 220:
                bullets.append(sent)
                if len(bullets) >= max_bullets:
                    break
    return bullets[:max_bullets]


# ---------------------------------------------------------------------------
# Reusable UI components
# ---------------------------------------------------------------------------

def _hero_box(slide, label: str, value: str, l, t, w, h) -> None:
    """Metric callout card: small label on top, large value below."""
    _rect(slide, l, t, w, h, _NAVY)
    mid = h * 0.42
    _label(slide, label,
           l + Inches(0.1), t + Inches(0.1),
           w - Inches(0.2), mid,
           size=10, color=_ICE_BLUE, align=PP_ALIGN.CENTER)
    _label(slide, value,
           l + Inches(0.1), t + mid,
           w - Inches(0.2), h - mid - Inches(0.1),
           bold=True, size=20, color=_WHITE, align=PP_ALIGN.CENTER)


def _kv_table(slide, header: str, rows: List[Tuple[str, str, Optional[int]]],
              lx, y, col_w, val_w, header_color=_NAVY) -> None:
    """Labelled two-column metric table with alternating row shading."""
    row_h = Inches(0.36)
    _label(slide, header, lx, y, col_w + val_w, Inches(0.28),
           bold=True, size=11, color=header_color)
    _rect(slide, lx, y + Inches(0.28), col_w + val_w, Inches(0.03), _ORANGE)
    for i, (lbl, val, cite) in enumerate(rows):
        ry = y + Inches(0.35) + i * row_h
        bg = _LIGHT_GREY if i % 2 == 0 else _WHITE
        _rect(slide, lx, ry, col_w + val_w, row_h - Inches(0.02), bg)
        _label(slide, lbl,
               lx + Inches(0.08), ry + Inches(0.05),
               col_w - Inches(0.1), row_h - Inches(0.08),
               size=10, color=_DARK_GREY)
        cite_suffix = f"  p.{cite}" if cite else ""
        _label(slide, val + cite_suffix,
               lx + col_w + Inches(0.05), ry + Inches(0.05),
               val_w - Inches(0.08), row_h - Inches(0.08),
               bold=True, size=10, color=_BLACK, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------

def _slide_1(prs, co: Dict, ds, date_str: str) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _NAVY)

    company_name = co.get("company_name", "Company Name")

    _label(slide, "Investment Committee Presentation",
           _MARGIN, Inches(1.85), _W - _MARGIN * 2, Inches(0.65),
           bold=True, size=36, color=_WHITE, align=PP_ALIGN.CENTER)

    _label(slide, company_name,
           _MARGIN, Inches(2.65), _W - _MARGIN * 2, Inches(0.60),
           bold=True, size=28, color=_ICE_BLUE, align=PP_ALIGN.CENTER)

    # Orange accent rule
    _rect(slide, Inches(5.4), Inches(3.38), Inches(2.53), Inches(0.055), _ORANGE)

    _label(slide, f"{date_str}     |     Prepared by: Meridian AI",
           _MARGIN, Inches(3.56), _W - _MARGIN * 2, Inches(0.38),
           size=13, color=_ICE_BLUE, align=PP_ALIGN.CENTER, italic=True)

    if ds and hasattr(ds, "grade"):
        _label(slide, f"Deal Score: {ds.total_score:.0%}  ·  Grade {ds.grade}  ·  {ds.recommendation}",
               _MARGIN, Inches(4.1), _W - _MARGIN * 2, Inches(0.38),
               size=13, color=_ORANGE, align=PP_ALIGN.CENTER)

    _label(slide, "◆  MERIDIAN AI  ·  PE Deal Intelligence Platform",
           _MARGIN, _H - Inches(0.52), _W - _MARGIN * 2, Inches(0.38),
           size=10, color=_MED_GREY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 2 — Deal Overview
# ---------------------------------------------------------------------------

def _slide_2(prs, co: Dict, fin: Dict, ds, currency: str) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Deal Overview")
    _footer(slide, 2)

    rev    = fin.get("revenue", {})
    ebitda = fin.get("ebitda", {})

    # Left column — company metadata
    lx   = _MARGIN
    info_y = Inches(0.82)

    _label(slide, co.get("company_name", "N/A"),
           lx, info_y, Inches(6.5), Inches(0.48),
           bold=True, size=20, color=_BLACK)

    industry    = co.get("industry", "")
    sub_industry = co.get("sub_industry", "")
    sector_line = industry
    if sub_industry and sub_industry != industry:
        sector_line += f"  ·  {sub_industry}"
    _label(slide, sector_line, lx, info_y + Inches(0.46), Inches(6.5), Inches(0.30),
           size=11, color=_DARK_GREY)

    info_rows = [
        ("Headquarters:",  _fmt_str(co.get("headquarters"))),
        ("Business Model:", _fmt_str(co.get("business_model"))),
        ("Founded:",        _fmt_str(co.get("founding_year"))),
        ("Employees:",      _fmt_str(co.get("employees") or co.get("employee_count"))),
    ]
    lbl_w = Inches(1.8)
    for i, (lbl, val) in enumerate(info_rows):
        ry = info_y + Inches(0.82) + i * Inches(0.30)
        _label(slide, lbl, lx, ry, lbl_w, Inches(0.28), bold=True, size=11, color=_DARK_GREY)
        _label(slide, val, lx + lbl_w + Inches(0.05), ry, Inches(4.2), Inches(0.28),
               size=11, color=_BLACK)

    # Right column — grade badge
    gx = Inches(8.6)
    gw = Inches(4.2)
    grade = getattr(ds, "grade", "C") if ds else "N/A"
    grade_color = _GRADE_COLORS.get(grade, _ORANGE)
    _rect(slide, gx, info_y, gw, Inches(1.05), grade_color)
    _label(slide, f"Grade  {grade}  ·  {getattr(ds, 'total_score', 0):.0%}",
           gx + Inches(0.1), info_y + Inches(0.08), gw - Inches(0.2), Inches(0.55),
           bold=True, size=26, color=_WHITE, align=PP_ALIGN.CENTER)
    _label(slide, getattr(ds, "recommendation", ""),
           gx + Inches(0.1), info_y + Inches(0.66), gw - Inches(0.2), Inches(0.32),
           size=12, color=_WHITE, align=PP_ALIGN.CENTER)

    # Hero metric boxes
    rev_ltm    = _with_currency(_fmt_num(rev.get("ltm")),              currency)
    ebitda_ltm = _with_currency(_fmt_num(ebitda.get("ltm")),           currency)
    margin     = _fmt_pct(ebitda.get("margin_ltm"))
    cagr       = _fmt_pct(rev.get("cagr_3yr"))
    boxes = [
        ("Revenue (LTM)",   rev_ltm),
        ("EBITDA (LTM)",    ebitda_ltm),
        ("EBITDA Margin",   margin),
        ("Revenue CAGR",    cagr),
    ]
    box_y = Inches(2.55)
    box_h = Inches(1.45)
    box_w = Inches(2.88)
    gap   = Inches(0.38)
    for i, (lbl, val) in enumerate(boxes):
        bx = _MARGIN + i * (box_w + gap)
        _hero_box(slide, lbl, val, bx, box_y, box_w, box_h)

    # Summary line
    summary = getattr(ds, "summary", "") if ds else ""
    if summary:
        _label(slide, summary[:180],
               _MARGIN, Inches(4.22), _W - _MARGIN * 2, Inches(0.48),
               size=11, color=_DARK_GREY, italic=True)

    # Deal context footnote
    deal = fin.get("deal_context") or {}  # not always in fin dict
    if not deal:
        from typing import Any as _Any
    reason = _fmt_str(co.get("reason_for_sale", ""))
    advisor = ""
    # Try extracted deal_context at top level
    # (handled gracefully — missing keys return "N/A")


# ---------------------------------------------------------------------------
# Slide 3 — Executive Summary
# ---------------------------------------------------------------------------

def _slide_3(prs, sections: List[Dict], co: Dict) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Executive Summary")
    _footer(slide, 3)

    raw_content = _clean_markdown(_find_section(sections, "executive summary", "executive"))
    _co_name_key = re.sub(r'[\s,\.]+', ' ', co.get("company_name", "")).strip().lower()

    def _is_company_name(text: str) -> bool:
        return re.sub(r'[\s,\.]+', ' ', text).strip().lower() == _co_name_key

    bullets: List[str] = []
    if raw_content:
        for para in raw_content.split("\n\n"):
            para = re.sub(r'\s*\n\s*', ' ', para).strip()
            if not para or len(para) < 20:
                continue
            if _is_company_name(para):
                continue
            # Truncate long paragraphs to first sentence
            if len(para) > 200:
                m = re.search(r'^(.+?[.!?])\s', para)
                para = m.group(1) if m else para[:200]
            bullets.append(para)
            if len(bullets) >= 6:
                break
        # If paragraph splitting yielded too few, fall back to sentence splitting
        if len(bullets) < 3:
            bullets = []
            full_flat = re.sub(r'\s*\n\s*', ' ', raw_content)
            for sent in re.split(r'(?<=[.!?])\s+', full_flat):
                sent = sent.strip()
                if len(sent) < 20 or _is_company_name(sent):
                    continue
                if sent not in bullets:
                    bullets.append(sent)
                if len(bullets) >= 6:
                    break
    if not bullets:
        bullets = _to_bullets(co.get("description", ""), 4)
    if not bullets:
        bullets = ["No executive summary available from memo."]

    body_y = Inches(0.82)
    body_h = _H - body_y - _FOOTER_H - Inches(0.2)
    lines = []
    for i, b in enumerate(bullets):
        lines.append((f"  •  {b}", {
            "size": 14, "color": _DARK_GREY,
            "space_before": 10 if i > 0 else 0,
        }))
    _multiline(slide, lines, _MARGIN, body_y, _W - _MARGIN * 2, body_h)


# ---------------------------------------------------------------------------
# Slide 4 — Company Overview
# ---------------------------------------------------------------------------

def _slide_4(prs, sections: List[Dict], co: Dict, customers: Dict) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Company Overview")
    _footer(slide, 4)

    raw_ov  = _clean_markdown(_find_section(sections, "company overview", "overview"))
    _co_nm  = re.sub(r'[\s,\.]+', ' ', co.get("company_name", "")).strip().lower()
    bullets: List[str] = []
    if raw_ov:
        for para in raw_ov.split("\n\n"):
            para = re.sub(r'\s*\n\s*', ' ', para).strip()
            if not para or len(para) < 20:
                continue
            if re.sub(r'[\s,\.]+', ' ', para).strip().lower() == _co_nm:
                continue
            if len(para) > 180:
                m = re.search(r'^(.+?[.!?])\s', para)
                para = m.group(1) if m else para[:180]
            bullets.append(para)
            if len(bullets) >= 5:
                break
        if len(bullets) < 2:
            bullets = _to_bullets(raw_ov, 5)
    if not bullets:
        bullets = _to_bullets(co.get("description", ""), 4)

    body_y = Inches(0.82)
    body_w = Inches(8.0)
    body_h = _H - body_y - _FOOTER_H - Inches(0.2)

    if bullets:
        lines = [(f"  •  {b}", {"size": 13, "color": _DARK_GREY,
                                 "space_before": 8 if i > 0 else 0})
                 for i, b in enumerate(bullets)]
        _multiline(slide, lines, _MARGIN, body_y, body_w, body_h)

    # Right — key facts card
    rx = Inches(8.7)
    rw = Inches(4.1)
    _label(slide, "KEY FACTS", rx, body_y, rw, Inches(0.28),
           bold=True, size=10, color=_NAVY)
    _rect(slide, rx, body_y + Inches(0.26), rw, Inches(0.03), _ORANGE)

    facts = [
        ("Industry",        co.get("industry", "N/A")),
        ("HQ",              co.get("headquarters", "N/A")),
        ("Founded",         _fmt_str(co.get("founding_year"))),
        ("Employees",       _fmt_str(co.get("employee_count") or co.get("employees"))),
        ("Biz Model",       co.get("business_model", "N/A")),
        ("Customers",       _fmt_str(customers.get("total_customers"))),
        ("Retention",       _fmt_pct(customers.get("customer_retention"))),
        ("Top Cust. Conc.", _fmt_pct(customers.get("top_customer_concentration"))),
    ]
    row_h = Inches(0.43)
    for i, (lbl, val) in enumerate(facts):
        fy = body_y + Inches(0.35) + i * row_h
        bg = _LIGHT_GREY if i % 2 == 0 else _WHITE
        _rect(slide, rx, fy, rw, row_h - Inches(0.02), bg)
        _label(slide, lbl, rx + Inches(0.08), fy + Inches(0.07),
               Inches(1.55), row_h - Inches(0.1), bold=True, size=9.5, color=_DARK_GREY)
        _label(slide, str(val)[:35], rx + Inches(1.68), fy + Inches(0.07),
               rw - Inches(1.78), row_h - Inches(0.1), size=9.5, color=_BLACK)


# ---------------------------------------------------------------------------
# Slide 5 — Financial Summary
# ---------------------------------------------------------------------------

def _slide_5(prs, fin: Dict, currency: str) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Financial Summary")
    _footer(slide, 5)

    rev    = fin.get("revenue", {})
    ebitda = fin.get("ebitda", {})

    cur_note = f"All figures in {currency}" if currency not in ("USD", "N/A", "") else "All figures in USD"
    _label(slide, cur_note,
           _MARGIN, Inches(0.71), Inches(6), Inches(0.20),
           size=8.5, color=_ORANGE, italic=True)

    def _n(v): return _with_currency(_fmt_num(v), currency)

    col_y  = Inches(0.95)
    col_w  = Inches(3.2)
    val_w  = Inches(2.6)
    gap    = Inches(0.55)

    left_rows: List[Tuple[str, str, Optional[int]]] = [
        ("Revenue (LTM)",       _n(rev.get("ltm")),                       _cite_page(rev.get("ltm_citation"))),
        ("Revenue (Prior Year)",_n(rev.get("prior_year")),                None),
        ("Revenue (2yr Ago)",   _n(rev.get("two_years_ago")),             None),
        ("CAGR (3yr)",          _fmt_pct(rev.get("cagr_3yr")),            _cite_page(rev.get("cagr_3yr_citation"))),
        ("EBITDA (LTM)",        _n(ebitda.get("ltm")),                    _cite_page(ebitda.get("ltm_citation"))),
        ("Adjusted EBITDA",     _n(ebitda.get("adjusted_ebitda_ltm")),    _cite_page(ebitda.get("adjusted_ebitda_ltm_citation"))),
        ("EBITDA Margin",       _fmt_pct(ebitda.get("margin_ltm")),       _cite_page(ebitda.get("margin_ltm_citation"))),
        ("Gross Margin",        _fmt_pct(fin.get("gross_margin")),        _cite_page(fin.get("gross_margin_citation"))),
        ("Net Income",          _n(fin.get("net_income")),                None),
    ]
    right_rows: List[Tuple[str, str, Optional[int]]] = [
        ("Total Debt",          _n(fin.get("debt")),                      None),
        ("Cash & Equivalents",  _n(fin.get("cash")),                      None),
        ("CapEx (Annual)",      _n(fin.get("capex")),                     _cite_page(fin.get("capex_citation"))),
        ("Free Cash Flow",      _n(fin.get("free_cash_flow")),            None),
        ("Recurring Rev %",     _fmt_pct(fin.get("recurring_revenue_pct")), None),
    ]
    # Add LTM source note as last right row
    ltm_src = _fmt_str(rev.get("ltm_source"))
    if ltm_src != "N/A":
        right_rows.append(("LTM Source", ltm_src[:38], None))

    _kv_table(slide, "INCOME STATEMENT", left_rows,
              _MARGIN, col_y, col_w, val_w)
    _kv_table(slide, "BALANCE SHEET & OTHER", right_rows,
              _MARGIN + col_w + val_w + gap, col_y, col_w, val_w)


# ---------------------------------------------------------------------------
# Slide 6 — Revenue Analysis
# ---------------------------------------------------------------------------

def _slide_6(prs, fin: Dict, currency: str) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Revenue Analysis")
    _footer(slide, 6)

    rev     = fin.get("revenue", {})
    cur_sym = currency if currency not in ("N/A", "") else "USD"
    history = rev.get("history") or []

    # Build chart series from history, falling back to simple 3-point series
    periods: List[str] = []
    vals_m:  List[float] = []

    if isinstance(history, list):
        for h in history:
            if not isinstance(h, dict):
                continue
            try:
                vals_m.append(float(h["value"]) / 1_000_000)
                periods.append(str(h.get("period", "")))
            except (KeyError, TypeError, ValueError):
                pass

    if not vals_m:
        for lbl, key in [("2yr Ago", "two_years_ago"), ("Prior Year", "prior_year"), ("LTM", "ltm")]:
            v = rev.get(key)
            if v is not None and str(v).lower() not in ("not_provided", "none"):
                try:
                    vals_m.append(float(v) / 1_000_000)
                    periods.append(lbl)
                except (TypeError, ValueError):
                    pass

    chart_y = Inches(0.82)
    chart_h = Inches(3.85)
    chart_w = Inches(7.6)

    if len(vals_m) >= 2:
        cd = ChartData()
        cd.categories = periods
        cd.add_series(f"Revenue ({cur_sym}M)", tuple(vals_m))

        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            _MARGIN, chart_y, chart_w, chart_h, cd,
        )
        chart = gf.chart
        chart.has_legend = False
        chart.has_title  = False
        try:
            for series in chart.series:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor(0x6F, 0x91, 0xFF)
        except Exception:
            pass
        try:
            chart.category_axis.tick_labels.font.size = Pt(9)
            chart.value_axis.has_major_gridlines = True
        except Exception:
            pass

        # YoY growth callouts — right of chart
        gx = _MARGIN + chart_w + Inches(0.2)
        gw = _W - gx - _MARGIN
        _label(slide, "YoY GROWTH", gx, chart_y, gw, Inches(0.28),
               bold=True, size=10, color=_NAVY)
        gy_base = chart_y + Inches(0.32)
        shown = 0
        for h in history:
            if not isinstance(h, dict):
                continue
            yoy = h.get("yoy_growth")
            if yoy is None or str(yoy).lower() in ("not_provided", "none"):
                continue
            try:
                pct = float(yoy)
                period = str(h.get("period", ""))[:20]
                col = _GREEN if pct >= 0 else _RED
                gy = gy_base + shown * Inches(0.40)
                _label(slide, f"{period}:", gx, gy, Inches(2.4), Inches(0.36),
                       size=9.5, color=_DARK_GREY)
                _label(slide, f"{pct:+.1%}", gx + Inches(2.5), gy, Inches(1.5), Inches(0.36),
                       bold=True, size=9.5, color=col)
                shown += 1
                if shown >= 8:
                    break
            except (TypeError, ValueError):
                pass
    else:
        _label(slide, "Insufficient revenue data for chart.",
               _MARGIN, chart_y + Inches(1.6), Inches(8), Inches(0.5),
               size=13, color=_MED_GREY, italic=True)

    # Revenue segments table
    segments = fin.get("revenue_by_segment") or []
    if isinstance(segments, list) and segments:
        seg_y = chart_y + chart_h + Inches(0.1)
        headers = ["Segment", "Revenue", "% of Total", "YoY Growth"]
        col_ws  = [Inches(3.2), Inches(1.8), Inches(1.5), Inches(1.5)]
        row_h   = Inches(0.30)

        _label(slide, "REVENUE BY SEGMENT",
               _MARGIN, seg_y, Inches(9), Inches(0.24),
               bold=True, size=9, color=_NAVY)
        seg_y += Inches(0.26)

        cx = _MARGIN
        for hdr, cw in zip(headers, col_ws):
            _rect(slide, cx, seg_y, cw - Inches(0.02), row_h, _NAVY)
            _label(slide, hdr, cx + Inches(0.06), seg_y + Inches(0.05),
                   cw - Inches(0.1), row_h - Inches(0.08),
                   bold=True, size=8.5, color=_WHITE)
            cx += cw

        total_w = sum(col_ws)
        for ri, seg in enumerate(segments[:5]):
            if not isinstance(seg, dict):
                continue
            ry = seg_y + row_h + ri * row_h
            bg = _LIGHT_GREY if ri % 2 == 0 else _WHITE
            _rect(slide, _MARGIN, ry, total_w, row_h - Inches(0.02), bg)
            vals = [
                str(seg.get("segment", ""))[:32],
                _with_currency(_fmt_num(seg.get("revenue")), currency),
                _fmt_pct(seg.get("pct_of_total")),
                _fmt_pct(seg.get("growth_rate")),
            ]
            cx = _MARGIN
            for val, cw in zip(vals, col_ws):
                _label(slide, val, cx + Inches(0.06), ry + Inches(0.04),
                       cw - Inches(0.1), row_h - Inches(0.06),
                       size=8.5, color=_BLACK)
                cx += cw


# ---------------------------------------------------------------------------
# Slide 7 — Growth Thesis
# ---------------------------------------------------------------------------

def _slide_7(prs, sections: List[Dict], growth_thesis: Dict) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Growth Thesis")
    _footer(slide, 7)

    raw_gt = _clean_markdown(_find_section(sections, "growth thesis", "growth"))
    bullets: List[str] = []

    # Paragraph-split first; each paragraph = one lever
    if raw_gt:
        for para in raw_gt.split("\n\n"):
            para = re.sub(r'\s*\n\s*', ' ', para).strip()
            if len(para) < 20:
                continue
            bullets.append(_truncate_at_word(para, 150))
            if len(bullets) >= 6:
                break
        # If paragraphs didn't yield enough, try line-by-line
        if len(bullets) < 3:
            bullets = [_truncate_at_word(b, 150) for b in _to_bullets(raw_gt, 6)]

    # Fallback to extracted growth_thesis dict fields
    if not bullets and isinstance(growth_thesis, dict):
        for lbl, key in [("Organic", "organic_levers"), ("M&A", "ma_opportunity"),
                         ("Expansion", "expansion_plans"), ("Technology", "technology_initiatives")]:
            items = growth_thesis.get(key, [])
            if isinstance(items, list):
                for it in items:
                    if it and str(it).lower() not in ("not_provided", "none"):
                        bullets.append(_truncate_at_word(str(it), 150))
            elif items and str(items).lower() not in ("not_provided", "none"):
                bullets.append(_truncate_at_word(f"{lbl}: {str(items)}", 150))
            if len(bullets) >= 6:
                break

    if not bullets:
        bullets = ["No growth thesis available from memo or extraction."]

    body_y = Inches(0.82)
    body_h = _H - body_y - _FOOTER_H - Inches(0.2)
    lines  = [(f"  •  {b}", {"size": 14, "color": _DARK_GREY,
                              "space_before": 10 if i > 0 else 0})
              for i, b in enumerate(bullets)]
    _multiline(slide, lines, _MARGIN, body_y, _W - _MARGIN * 2, body_h)


# ---------------------------------------------------------------------------
# Slide 8 — Key Risks & Mitigants
# ---------------------------------------------------------------------------

def _slide_8(prs, risks, sections: List[Dict]) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Key Risks & Mitigants")
    _footer(slide, 8)

    # Build rows from automated RiskFlag objects
    rows: List[Tuple[str, str, str]] = []
    for r in (risks or []):
        title = getattr(r, "title", "") or getattr(r, "description", "")[:60] or ""
        sev   = getattr(r, "severity", "Medium") or "Medium"
        # Field is 'mitigant', not 'mitigation'
        mit   = (getattr(r, "mitigant", None) or getattr(r, "mitigation", None) or "").strip()
        rows.append((title[:90], sev, (mit[:100] if mit else "Under diligence")))

    # Supplement with memo qualitative risks (with actual mitigant text)
    if len(rows) < 8:
        memo_content = _clean_markdown(_find_section(sections, "key risks", "risk", "mitigant"))
        # Try to parse numbered blocks with mitigants
        blocks = re.split(r'\n(?=\d+[\.\)])', memo_content)
        for block in blocks:
            block = block.strip()
            if not block or not re.match(r'^\d+', block):
                continue
            # Split on "Mitigant:" label
            parts = re.split(r'\n?\s*(?:\*\*?)?Mitigant(?:\*\*?)?\s*:', block, maxsplit=1, flags=re.IGNORECASE)
            risk_text = re.sub(r'^\d+[\.\)]\s*', '', parts[0]).replace('**', '').strip()
            mitigant  = parts[1].strip().replace('**', '') if len(parts) > 1 else "Under diligence"
            # First line of risk_text = title
            risk_title = risk_text.split('\n')[0].strip()[:90]
            if not risk_title or len(risk_title) < 10:
                continue
            if any(risk_title[:28] in r[0] for r in rows):
                continue
            rows.append((risk_title, "Medium", mitigant[:100]))
            if len(rows) >= 12:
                break

    if not rows:
        rows = [("No automated risks flagged — perform manual risk review.", "Low", "N/A")]

    table_y = Inches(0.80)
    row_h   = Inches(0.36)
    c_risk  = Inches(5.6)
    c_sev   = Inches(1.2)
    c_mit   = _W - _MARGIN * 2 - c_risk - c_sev
    headers = [("Risk / Observation", c_risk), ("Severity", c_sev), ("Mitigant / Diligence Focus", c_mit)]

    cx = _MARGIN
    for hdr, cw in headers:
        _rect(slide, cx, table_y, cw - Inches(0.02), Inches(0.33), _NAVY)
        _label(slide, hdr, cx + Inches(0.08), table_y + Inches(0.06),
               cw - Inches(0.14), Inches(0.24), bold=True, size=9.5, color=_WHITE)
        cx += cw

    max_rows = 14
    for ri, (risk_title, severity, mitigation) in enumerate(rows[:max_rows]):
        ry  = table_y + Inches(0.33) + ri * row_h
        bg  = _LIGHT_GREY if ri % 2 == 0 else _WHITE
        sc  = _SEV_COLORS.get(severity, _AMBER)

        _rect(slide, _MARGIN, ry, c_risk + c_sev + c_mit, row_h - Inches(0.02), bg)

        _label(slide, risk_title,
               _MARGIN + Inches(0.08), ry + Inches(0.05),
               c_risk - Inches(0.14), row_h - Inches(0.08),
               size=9.5, color=_BLACK)

        badge_lx = _MARGIN + c_risk + Inches(0.10)
        badge_w  = c_sev - Inches(0.2)
        _rect(slide, badge_lx, ry + Inches(0.07), badge_w, row_h - Inches(0.16), sc)
        _label(slide, severity,
               badge_lx, ry + Inches(0.08), badge_w, row_h - Inches(0.18),
               bold=True, size=7.5, color=_WHITE, align=PP_ALIGN.CENTER)

        _label(slide, mitigation,
               _MARGIN + c_risk + c_sev + Inches(0.08), ry + Inches(0.05),
               c_mit - Inches(0.14), row_h - Inches(0.08),
               size=9.5, color=_DARK_GREY)


# ---------------------------------------------------------------------------
# Slide 9 — Comparable Companies
# ---------------------------------------------------------------------------

def _slide_9(prs, comps, market: Dict, extracted_data: Optional[Dict] = None) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Comparable Companies")
    _footer(slide, 9)

    _label(slide,
           "Note: Financial multiples (EV/EBITDA, EV/Revenue) require "
           "PitchBook / Capital IQ API integration.",
           _MARGIN, Inches(0.74), _W - _MARGIN * 2, Inches(0.22),
           size=9, color=_ORANGE, italic=True)

    cp = market.get("competitive_position", "")
    if cp and cp not in ("not_provided", "N/A"):
        _label(slide, f"Competitive Position: {cp}",
               _MARGIN, Inches(1.0), Inches(8), Inches(0.28),
               bold=True, size=12, color=_NAVY)

    # Collect identified peers — try multiple possible field names
    _raw_peers = (
        market.get("key_competitors")
        or market.get("competitors")
        or market.get("direct_competitors")
        or (extracted_data or {}).get("key_competitors")
        or []
    )
    key_competitors = []
    for c in (_raw_peers or []):
        name = c.get("name", str(c)) if isinstance(c, dict) else str(c)
        name = name.strip()
        if name and name.lower() not in ("not_provided", "none", "n/a"):
            key_competitors.append(name)

    table_y = Inches(1.36)

    if comps:
        col_ws  = [Inches(3.8), Inches(1.9), Inches(1.9), Inches(4.9)]
        headers = ["Company", "EV/EBITDA", "EV/Revenue", "Rationale"]
        row_h   = Inches(0.36)

        cx = _MARGIN
        for hdr, cw in zip(headers, col_ws):
            _rect(slide, cx, table_y, cw - Inches(0.02), Inches(0.32), _NAVY)
            _label(slide, hdr, cx + Inches(0.08), table_y + Inches(0.06),
                   cw - Inches(0.14), Inches(0.24), bold=True, size=9.5, color=_WHITE)
            cx += cw

        for ri, comp in enumerate(comps[:12]):
            ry  = table_y + Inches(0.32) + ri * row_h
            bg  = _LIGHT_GREY if ri % 2 == 0 else _WHITE
            _rect(slide, _MARGIN, ry, sum(col_ws), row_h - Inches(0.02), bg)
            ev_ebitda = getattr(comp, "ev_ebitda", None)
            ev_rev    = getattr(comp, "ev_revenue", None)
            vals = [
                str(getattr(comp, "name", "N/A"))[:45],
                f"{ev_ebitda:.1f}x" if ev_ebitda else "N/A",
                f"{ev_rev:.1f}x"    if ev_rev    else "N/A",
                str(getattr(comp, "rationale", ""))[:80],
            ]
            cx = _MARGIN
            for val, cw in zip(vals, col_ws):
                _label(slide, val, cx + Inches(0.08), ry + Inches(0.05),
                       cw - Inches(0.14), row_h - Inches(0.08), size=9.5, color=_BLACK)
                cx += cw

        # Identified peers row below comps table
        if key_competitors:
            peers_y = table_y + Inches(0.32) + min(len(comps), 12) * row_h + Inches(0.22)
            peer_str = "Identified Peers:  " + "   |   ".join(key_competitors[:8])
            _label(slide, peer_str, _MARGIN, peers_y, _W - _MARGIN * 2, Inches(0.32),
                   size=9.5, color=_DARK_GREY, italic=True)
    else:
        # No LLM comps — show identified peers prominently
        if key_competitors:
            _label(slide, "Identified Peers from Document Analysis",
                   _MARGIN, table_y, _W - _MARGIN * 2, Inches(0.30),
                   bold=True, size=13, color=_NAVY)
            _rect(slide, _MARGIN, table_y + Inches(0.30), Inches(3.0), Inches(0.04), _ORANGE)
            lines = [(f"  •  {c}", {"size": 12, "color": _DARK_GREY,
                                     "space_before": 6 if i > 0 else 0})
                     for i, c in enumerate(key_competitors[:10])]
            _multiline(slide, lines, _MARGIN, table_y + Inches(0.42),
                       _W - _MARGIN * 2, Inches(4.0))
            _label(slide,
                   "No transaction multiples available — financial comps require PitchBook / Capital IQ.",
                   _MARGIN, table_y + Inches(4.6), _W - _MARGIN * 2, Inches(0.26),
                   size=9, color=_MED_GREY, italic=True)
        else:
            _label(slide,
                   "No comparable companies or identified peers found in this document.\n"
                   "Comparable analysis requires PitchBook / Capital IQ API for sector transaction multiples.",
                   _MARGIN, table_y, _W - _MARGIN * 2, Inches(0.65),
                   size=12, color=_DARK_GREY)


# ---------------------------------------------------------------------------
# Slide 10 — Deal Scoring
# ---------------------------------------------------------------------------

def _slide_10(prs, ds) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _OFF_WHITE)
    _header(slide, "Deal Scoring")
    _footer(slide, 10)

    if not ds:
        _label(slide, "Deal score not available.",
               _MARGIN, Inches(2), Inches(8), Inches(0.5),
               size=14, color=_MED_GREY)
        return

    # Grade badge — top right
    grade = getattr(ds, "grade", "C")
    gc = _GRADE_COLORS.get(grade, _ORANGE)
    gx = _W - _MARGIN - Inches(3.6)
    _rect(slide, gx, Inches(0.82), Inches(3.6), Inches(1.15), gc)
    _label(slide, f"Grade  {grade}",
           gx + Inches(0.1), Inches(0.88), Inches(3.4), Inches(0.58),
           bold=True, size=32, color=_WHITE, align=PP_ALIGN.CENTER)
    _label(slide, f"{ds.total_score:.0%}  Overall Score",
           gx + Inches(0.1), Inches(1.48), Inches(3.4), Inches(0.32),
           size=13, color=_WHITE, align=PP_ALIGN.CENTER)

    _label(slide, getattr(ds, "recommendation", ""),
           _MARGIN, Inches(0.82), Inches(8.6), Inches(0.45),
           bold=True, size=18, color=_NAVY)

    # Dimension bars
    dims    = getattr(ds, "dimensions", [])
    bar_y   = Inches(1.55)
    lbl_w   = Inches(2.9)
    bar_max = Inches(5.6)
    bar_h   = Inches(0.38)

    for i, dim in enumerate(dims):
        dy       = bar_y + i * Inches(0.74)
        name     = getattr(dim, "dimension", "")
        score    = float(getattr(dim, "score", 0))
        weight   = float(getattr(dim, "weight", 0))
        rationale = str(getattr(dim, "rationale", ""))

        _label(slide, f"{name}  ({weight:.0%})",
               _MARGIN, dy, lbl_w, bar_h,
               bold=True, size=11, color=_DARK_GREY)

        # Bar background
        _rect(slide, _MARGIN + lbl_w + Inches(0.1), dy + Inches(0.06),
              bar_max, bar_h - Inches(0.12), _LIGHT_GREY)
        # Bar fill — color by score tier
        fill_c = _GREEN if score >= 0.70 else (_AMBER if score >= 0.50 else _RED)
        fill_w = max(Inches(0.04), bar_max * score)
        _rect(slide, _MARGIN + lbl_w + Inches(0.1), dy + Inches(0.06),
              fill_w, bar_h - Inches(0.12), fill_c)
        # Percentage
        _label(slide, f"{score:.0%}",
               _MARGIN + lbl_w + Inches(0.1) + bar_max + Inches(0.1),
               dy + Inches(0.06), Inches(0.65), bar_h - Inches(0.12),
               bold=True, size=11, color=_BLACK)
        # Rationale
        if rationale:
            _label(slide, rationale[:120],
                   _MARGIN + lbl_w + Inches(0.1), dy + bar_h - Inches(0.04),
                   bar_max + Inches(0.8), Inches(0.28),
                   size=8, color=_MED_GREY, italic=True)

    summary = str(getattr(ds, "summary", ""))
    if summary:
        sy = bar_y + len(dims) * Inches(0.74) + Inches(0.1)
        _label(slide, summary[:220],
               _MARGIN, sy, _W - _MARGIN * 2 - Inches(0.2), Inches(0.55),
               size=10, color=_DARK_GREY, italic=True)


# ---------------------------------------------------------------------------
# Slide 11 — Key Diligence Questions + Recommendation
# ---------------------------------------------------------------------------

def _slide_11(prs, sections: List[Dict], ds) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, _NAVY)
    _footer(slide, 11)

    _label(slide, "Key Diligence Questions",
           _MARGIN, Inches(0.28), _W - _MARGIN * 2, Inches(0.58),
           bold=True, size=28, color=_WHITE)
    _rect(slide, _MARGIN, Inches(0.92), Inches(3.2), Inches(0.055), _ORANGE)

    content   = _find_section(sections, "diligence", "due diligence",
                               "key questions", "key diligence")
    # Prefer explicitly numbered/bulleted items to avoid preamble prose
    questions = _parse_numbered_items(content, 7) or _to_bullets(content, 7)
    if not questions:
        questions = [
            "Validate revenue quality: verify recurring vs. one-time revenue split.",
            "Assess customer concentration: obtain full customer list with revenue by account.",
            "Management depth: evaluate bench strength beyond founder / CEO.",
            "Confirm financials with audited statements and audit opinion.",
            "Model downside scenarios: 20% revenue decline impact on covenant compliance.",
        ]

    body_y = Inches(1.05)
    body_h = _H - body_y - _FOOTER_H - Inches(0.85)
    lines  = [(f"  {i+1}.  {q}", {"size": 13, "color": _ICE_BLUE,
                                    "space_before": 9 if i > 0 else 0})
              for i, q in enumerate(questions)]
    _multiline(slide, lines, _MARGIN, body_y, _W - _MARGIN * 2, body_h)

    # Recommendation banner at bottom
    rec = getattr(ds, "recommendation", "") if ds else ""
    if rec:
        rec_y = _H - _FOOTER_H - Inches(0.62)
        _rect(slide, _MARGIN, rec_y, _W - _MARGIN * 2, Inches(0.52), _ORANGE)
        _label(slide, f"Recommendation: {rec}",
               _MARGIN + Inches(0.2), rec_y + Inches(0.06),
               _W - _MARGIN * 2 - Inches(0.4), Inches(0.42),
               bold=True, size=16, color=_WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def generate_ic_deck(result: "AnalysisResult", output_path: str) -> None:
    """Generate an 11-slide IC deck PowerPoint from an AnalysisResult.

    Args:
        result:      Pipeline AnalysisResult (extracted_data, memo, risks, comps, deal_score).
        output_path: Destination .pptx file path (created if necessary).
    """
    prs = _new_prs()

    ed        = result.extracted_data
    co        = ed.get("company_overview", {})
    fin       = ed.get("financials", {})
    customers = ed.get("customers", {})
    market    = ed.get("market", {})
    growth    = ed.get("growth_thesis", {})
    ds        = result.deal_score
    risks     = result.risks or []
    comps     = result.comps or []

    currency = str(fin.get("currency") or "USD")
    if currency.lower() in ("none", "not_provided", ""):
        currency = "USD"

    memo_data = _parse_memo(result.memo)
    sections  = memo_data.get("sections", [])
    date_str  = (memo_data.get("date") or datetime.today().strftime("%B %d, %Y"))

    _slide_1(prs, co, ds, date_str)
    _slide_2(prs, co, fin, ds, currency)
    _slide_3(prs, sections, co)
    _slide_4(prs, sections, co, customers)
    _slide_5(prs, fin, currency)
    _slide_6(prs, fin, currency)
    _slide_7(prs, sections, growth)
    _slide_8(prs, risks, sections)
    _slide_9(prs, comps, market, extracted_data=ed)
    _slide_10(prs, ds)
    _slide_11(prs, sections, ds)

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
